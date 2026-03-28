"""
pptxプレゼンテーション生成スクリプト
Generate PowerPoint presentations (English + Japanese) with 1 figure/table per slide.

- Code-generated figures: embedded as images (not editable)
- Flow diagrams and concept diagrams: built as editable PowerPoint shapes
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Inches, Pt

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "output")
PAPER_DIR = os.path.join(SCRIPT_DIR, "papers")
os.makedirs(PAPER_DIR, exist_ok=True)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# Color palette
COLOR_TITLE_BG = RGBColor(0, 51, 102)
COLOR_TITLE_TEXT = RGBColor(255, 255, 255)
COLOR_HEADING = RGBColor(0, 51, 102)
COLOR_BODY = RGBColor(51, 51, 51)
COLOR_ACCENT1 = RGBColor(0, 102, 153)
COLOR_ACCENT2 = RGBColor(204, 85, 0)
COLOR_LIGHT_BG = RGBColor(240, 248, 255)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(180, 180, 180)
COLOR_GREEN = RGBColor(46, 139, 87)
COLOR_RED = RGBColor(178, 34, 34)
COLOR_ORANGE = RGBColor(210, 105, 30)
COLOR_BLUE = RGBColor(70, 130, 180)
COLOR_PURPLE = RGBColor(128, 0, 128)
COLOR_TEAL = RGBColor(0, 128, 128)

SPECIES_JA = {
    "Sperm_Whale": "マッコウクジラ",
    "Humpback_Whale": "ザトウクジラ",
    "Killer_Whale": "シャチ",
    "Fin_Finback_Whale": "ナガスクジラ",
    "Bottlenose_Dolphin": "バンドウイルカ",
    "Beluga_White_Whale": "シロイルカ",
}

SPECIES_EN = {
    "Sperm_Whale": "Sperm Whale",
    "Humpback_Whale": "Humpback Whale",
    "Killer_Whale": "Killer Whale",
    "Fin_Finback_Whale": "Fin Whale",
    "Bottlenose_Dolphin": "Bottlenose Dolphin",
    "Beluga_White_Whale": "Beluga Whale",
}

TARGET_SPECIES = [
    "Sperm_Whale",
    "Humpback_Whale",
    "Killer_Whale",
    "Fin_Finback_Whale",
    "Bottlenose_Dolphin",
    "Beluga_White_Whale",
]

CLICK_SPECIES = ["Sperm_Whale", "Killer_Whale", "Bottlenose_Dolphin"]


# ============================================================
# Helper functions
# ============================================================


def add_shape(slide, left, top, width, height, shape_type=MSO_SHAPE.RECTANGLE,
              fill_color=None, line_color=None, line_width=Pt(1)):
    """Add a shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=Pt(12), font_color=COLOR_BODY, bold=False,
             alignment=PP_ALIGN.CENTER, font_name=None):
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return tf


def add_text_to_frame(tf, text, font_size=Pt(12), font_color=COLOR_BODY, bold=False,
                      alignment=PP_ALIGN.CENTER, font_name=None):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return p


def add_arrow_connector(slide, start_left, start_top, end_left, end_top, color=COLOR_GRAY, width=Pt(2)):
    """Add an arrow shape between two points."""
    # Use a line shape (freeform connector approximation via thin rectangle + arrowhead)
    # We'll use a right arrow shape for horizontal or down arrow for vertical
    dx = end_left - start_left
    dy = end_top - start_top
    if abs(dy) > abs(dx):
        # Vertical arrow (down)
        arrow_left = start_left - Cm(0.3)
        arrow_top = start_top
        arrow_width = Cm(0.6)
        arrow_height = dy
        shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow_left, arrow_top, arrow_width, arrow_height)
    else:
        # Horizontal arrow (right)
        arrow_left = start_left
        arrow_top = start_top - Cm(0.3)
        arrow_width = dx
        arrow_height = Cm(0.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, arrow_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, prs, num, total):
    """Add slide number in bottom right."""
    txBox = slide.shapes.add_textbox(Cm(30), Cm(18.0), Cm(3), Cm(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY
    p.alignment = PP_ALIGN.RIGHT


def make_title_slide(prs, title, subtitle):
    """Create a title slide with colored background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Full background rectangle
    bg = add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, SLIDE_HEIGHT,
                   fill_color=COLOR_TITLE_BG)

    # Title
    txBox = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(28), Cm(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_TITLE_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Cm(3), Cm(11), Cm(28), Cm(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.alignment = PP_ALIGN.CENTER

    return slide


def make_section_slide(prs, section_title, section_num=None):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Accent bar at top
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.4), fill_color=COLOR_TITLE_BG)

    # Section title
    prefix = f"{section_num}. " if section_num else ""
    txBox = slide.shapes.add_textbox(Cm(3), Cm(6), Cm(28), Cm(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{prefix}{section_title}"
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    return slide


def make_image_slide(prs, title, image_path, caption=""):
    """Create a slide with a single code-output image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Accent bar
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    # Title
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Image (centered, large)
    if image_path and os.path.exists(image_path):
        # Calculate size to fit while maintaining aspect ratio
        max_width = Cm(30)
        max_height = Cm(14.5)
        left = Cm(2)
        top = Cm(2.3)
        slide.shapes.add_picture(image_path, left, top, width=max_width)

    # Caption
    if caption:
        txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(32), Cm(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_GRAY
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

    return slide


def make_table_slide(prs, title, headers, rows, caption=""):
    """Create a slide with an editable table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Accent bar
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    # Title
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Cm(30)
    tbl_height = Cm(n_rows * 1.2)
    left = Cm(2)
    top = Cm(2.5)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width, tbl_height)
    table = table_shape.table

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER
        # Header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE_BG

    # Data rows
    for r, row_data in enumerate(rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_BODY
                p.alignment = PP_ALIGN.CENTER
            # Alternating row colors
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    # Caption
    if caption:
        txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(32), Cm(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_GRAY
        p2.font.italic = True

    return slide


# ============================================================
# Editable flow/concept diagram slides
# ============================================================


def make_analysis_pipeline_flow(prs, lang="ja"):
    """Create editable analysis pipeline flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    title_text = "解析パイプライン概要" if lang == "ja" else "Analysis Pipeline Overview"
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    # Pipeline boxes
    if lang == "ja":
        steps = [
            ("データ取得", "Watkins WMMS\n1,357サンプル, 32種", COLOR_ACCENT1),
            ("前処理", "音声デコード\nモノラル変換\nリサンプリング (16kHz)", COLOR_BLUE),
            ("特徴抽出", "スペクトログラム\nICI検出\nバイスペクトル\nエントロピー\n時間構造", COLOR_GREEN),
            ("比較解析", "種間比較\nCDMA直交性検定\nMann-Whitney U検定", COLOR_ORANGE),
            ("結果", "論文レポート\n可視化 (58図)\n統計的検証", COLOR_PURPLE),
        ]
    else:
        steps = [
            ("Data Acquisition", "Watkins WMMS\n1,357 samples, 32 species", COLOR_ACCENT1),
            ("Preprocessing", "Audio decode\nStereo→Mono\nResample (16kHz)", COLOR_BLUE),
            ("Feature Extraction", "Spectrogram\nICI detection\nBispectrum\nEntropy\nTemporal", COLOR_GREEN),
            ("Comparative Analysis", "Cross-species\nCDMA orthogonality\nMann-Whitney U test", COLOR_ORANGE),
            ("Output", "Paper report\nVisualizations (58 figs)\nStatistical validation", COLOR_PURPLE),
        ]

    box_width = Cm(5.4)
    box_height = Cm(6.5)
    start_left = Cm(1.5)
    top = Cm(3.5)
    gap = Cm(1.1)

    for i, (title, desc, color) in enumerate(steps):
        left = start_left + i * (box_width + gap)

        # Box with rounded corners
        box = add_shape(slide, left, top, box_width, box_height,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=COLOR_WHITE, line_color=color, line_width=Pt(2))

        # Title bar at top of box
        title_bar = add_shape(slide, left, top, box_width, Cm(1.5),
                              shape_type=MSO_SHAPE.RECTANGLE, fill_color=color)
        set_text(title_bar, title, font_size=Pt(12), font_color=COLOR_WHITE, bold=True)

        # Description text
        desc_box = slide.shapes.add_textbox(left + Cm(0.3), top + Cm(1.8),
                                            box_width - Cm(0.6), box_height - Cm(2.2))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLOR_BODY
        dp.alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < len(steps) - 1:
            arrow_left = left + box_width
            arrow_top = top + box_height / 2
            arrow = add_shape(slide, arrow_left, arrow_top - Cm(0.25),
                              gap, Cm(0.5),
                              shape_type=MSO_SHAPE.RIGHT_ARROW,
                              fill_color=COLOR_GRAY)

    # Step numbers
    for i in range(len(steps)):
        left = start_left + i * (box_width + gap)
        num_shape = add_shape(slide, left + box_width / 2 - Cm(0.5), top - Cm(1.0),
                              Cm(1.0), Cm(1.0),
                              shape_type=MSO_SHAPE.OVAL,
                              fill_color=COLOR_TITLE_BG)
        set_text(num_shape, str(i + 1), font_size=Pt(14), font_color=COLOR_WHITE, bold=True)

    # Footer note
    note_text = "※ すべての要素は編集可能です" if lang == "ja" else "Note: All elements are editable"
    txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = note_text
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.italic = True

    return slide


def make_hypothesis_diagram(prs, lang="ja"):
    """Create editable hypothesis concept diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    title_text = "検証仮説の概念図" if lang == "ja" else "Conceptual Diagram of Hypotheses Under Test"
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    # --- Left side: CDMA hypothesis ---
    if lang == "ja":
        cdma_title = "仮説1: CDMA的符号分割"
        cdma_items = [
            ("発信者A", "固有スペクトル符号 A"),
            ("発信者B", "固有スペクトル符号 B"),
            ("発信者C", "固有スペクトル符号 C"),
        ]
        cdma_receiver = "受信者"
        cdma_decode = "符号Aのみ復号\n→ Aのメッセージ取得"
        cdma_result = "検証結果: 種内相関 > 種間相関\n(p = 1.57×10⁻³) → 支持"
    else:
        cdma_title = "Hypothesis 1: CDMA-like Code Division"
        cdma_items = [
            ("Sender A", "Unique spectral code A"),
            ("Sender B", "Unique spectral code B"),
            ("Sender C", "Unique spectral code C"),
        ]
        cdma_receiver = "Receiver"
        cdma_decode = "Decode only code A\n→ Extract A's message"
        cdma_result = "Result: Within > Between correlation\n(p = 1.57×10⁻³) → Supported"

    # CDMA section header
    cdma_header = add_shape(slide, Cm(0.5), Cm(2.5), Cm(16), Cm(1.2),
                            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                            fill_color=COLOR_ACCENT1)
    set_text(cdma_header, cdma_title, font_size=Pt(14), font_color=COLOR_WHITE, bold=True)

    # Sender boxes
    colors_sender = [COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE]
    for i, (sender, code) in enumerate(cdma_items):
        y = Cm(4.2) + i * Cm(2.0)
        box = add_shape(slide, Cm(1), y, Cm(4), Cm(1.5),
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=colors_sender[i], line_color=colors_sender[i])
        set_text(box, sender, font_size=Pt(11), font_color=COLOR_WHITE, bold=True)

        # Code label
        code_box = add_shape(slide, Cm(5.5), y, Cm(5), Cm(1.5),
                             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                             fill_color=COLOR_WHITE, line_color=colors_sender[i], line_width=Pt(1.5))
        set_text(code_box, code, font_size=Pt(9), font_color=colors_sender[i])

        # Arrow to channel
        add_shape(slide, Cm(10.5), y + Cm(0.5), Cm(1), Cm(0.4),
                  shape_type=MSO_SHAPE.RIGHT_ARROW, fill_color=COLOR_GRAY)

    # Channel/medium box
    channel_y = Cm(4.2)
    channel = add_shape(slide, Cm(11.5), channel_y, Cm(2.0), Cm(5.5),
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=RGBColor(230, 230, 250),
                        line_color=COLOR_PURPLE, line_width=Pt(1.5))
    ch_text = "海中\n伝搬" if lang == "ja" else "Ocean\nChannel"
    set_text(channel, ch_text, font_size=Pt(10), font_color=COLOR_PURPLE, bold=True)

    # Arrow to receiver
    add_shape(slide, Cm(13.5), Cm(6.5), Cm(1), Cm(0.4),
              shape_type=MSO_SHAPE.RIGHT_ARROW, fill_color=COLOR_GRAY)

    # Receiver
    recv = add_shape(slide, Cm(14.5), Cm(5.5), Cm(2.0), Cm(2.5),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                     fill_color=COLOR_ACCENT1, line_color=COLOR_ACCENT1)
    set_text(recv, cdma_receiver, font_size=Pt(11), font_color=COLOR_WHITE, bold=True)

    # Result box
    result_box = add_shape(slide, Cm(0.5), Cm(10.5), Cm(16), Cm(2.0),
                           shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                           fill_color=RGBColor(240, 255, 240),
                           line_color=COLOR_GREEN, line_width=Pt(2))
    set_text(result_box, cdma_result, font_size=Pt(11), font_color=COLOR_GREEN, bold=True)

    # --- Right side: Beat frequency hypothesis ---
    if lang == "ja":
        beat_title = "仮説2: うなり周波数効果"
        signal1 = "信号 f₁"
        signal2 = "信号 f₂\n(時間差 Δt)"
        interference = "干渉\n(非線形相互作用)"
        beat_comp = "うなり成分\nf₁ - f₂"
        beat_result = "検証結果: 全種でバイコヒーレンス > 0\nマッコウクジラが最大 → 部分的支持"
    else:
        beat_title = "Hypothesis 2: Beat Frequency Effect"
        signal1 = "Signal f₁"
        signal2 = "Signal f₂\n(time delay Δt)"
        interference = "Interference\n(nonlinear interaction)"
        beat_comp = "Beat component\nf₁ - f₂"
        beat_result = "Result: Bicoherence > 0 in all species\nSperm whale highest → Partial support"

    # Beat section header
    beat_header = add_shape(slide, Cm(17.5), Cm(2.5), Cm(16), Cm(1.2),
                            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                            fill_color=COLOR_ACCENT2)
    set_text(beat_header, beat_title, font_size=Pt(14), font_color=COLOR_WHITE, bold=True)

    # Signal 1
    s1 = add_shape(slide, Cm(18), Cm(4.5), Cm(5), Cm(1.5),
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                   fill_color=COLOR_BLUE, line_color=COLOR_BLUE)
    set_text(s1, signal1, font_size=Pt(11), font_color=COLOR_WHITE, bold=True)

    # Signal 2
    s2 = add_shape(slide, Cm(18), Cm(6.5), Cm(5), Cm(1.5),
                   shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                   fill_color=COLOR_RED, line_color=COLOR_RED)
    set_text(s2, signal2, font_size=Pt(10), font_color=COLOR_WHITE, bold=True)

    # Arrows to interference
    add_shape(slide, Cm(23), Cm(5.0), Cm(1.5), Cm(0.4),
              shape_type=MSO_SHAPE.RIGHT_ARROW, fill_color=COLOR_BLUE)
    add_shape(slide, Cm(23), Cm(7.0), Cm(1.5), Cm(0.4),
              shape_type=MSO_SHAPE.RIGHT_ARROW, fill_color=COLOR_RED)

    # Interference
    interf = add_shape(slide, Cm(24.5), Cm(4.8), Cm(4.5), Cm(3.0),
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                       fill_color=RGBColor(255, 245, 230),
                       line_color=COLOR_ACCENT2, line_width=Pt(2))
    set_text(interf, interference, font_size=Pt(10), font_color=COLOR_ACCENT2, bold=True)

    # Arrow down to beat
    add_shape(slide, Cm(26.5), Cm(7.8), Cm(0.5), Cm(1.2),
              shape_type=MSO_SHAPE.DOWN_ARROW, fill_color=COLOR_ACCENT2)

    # Beat component
    beat = add_shape(slide, Cm(24.5), Cm(9.0), Cm(4.5), Cm(1.5),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                     fill_color=COLOR_PURPLE, line_color=COLOR_PURPLE)
    set_text(beat, beat_comp, font_size=Pt(10), font_color=COLOR_WHITE, bold=True)

    # Result box
    result_box2 = add_shape(slide, Cm(17.5), Cm(10.5), Cm(16), Cm(2.0),
                            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                            fill_color=RGBColor(255, 255, 230),
                            line_color=COLOR_ORANGE, line_width=Pt(2))
    set_text(result_box2, beat_result, font_size=Pt(11), font_color=COLOR_ORANGE, bold=True)

    # Footer
    note_text = "※ すべての要素は編集可能です" if lang == "ja" else "Note: All elements are editable"
    txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = note_text
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.italic = True

    return slide


def make_species_overview_diagram(prs, lang="ja"):
    """Create editable species overview diagram with frequency ranges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    title_text = "解析対象6種の音響特性概要" if lang == "ja" else "Acoustic Characteristics of 6 Target Species"
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    if lang == "ja":
        species_info = [
            ("マッコウクジラ", "クリック列 (コーダ)", "2,309 Hz", "7.32 bits", COLOR_BLUE),
            ("ザトウクジラ", "歌 (Zipf分布)", "593 Hz", "4.47 bits", COLOR_GREEN),
            ("シャチ", "コール・ホイッスル", "1,804 Hz", "6.66 bits", COLOR_ORANGE),
            ("ナガスクジラ", "超低周波パルス", "38 Hz", "2.28 bits", COLOR_TEAL),
            ("バンドウイルカ", "ホイッスル・クリック", "4,021 Hz", "5.79 bits", COLOR_PURPLE),
            ("シロイルカ", "多様な発声", "1,251 Hz", "6.19 bits", COLOR_RED),
        ]
        headers_text = ["種名", "発声タイプ", "重心周波数", "エントロピー"]
    else:
        species_info = [
            ("Sperm Whale", "Click trains (codas)", "2,309 Hz", "7.32 bits", COLOR_BLUE),
            ("Humpback Whale", "Songs (Zipf distribution)", "593 Hz", "4.47 bits", COLOR_GREEN),
            ("Killer Whale", "Calls & whistles", "1,804 Hz", "6.66 bits", COLOR_ORANGE),
            ("Fin Whale", "Infrasonic pulses", "38 Hz", "2.28 bits", COLOR_TEAL),
            ("Bottlenose Dolphin", "Whistles & clicks", "4,021 Hz", "5.79 bits", COLOR_PURPLE),
            ("Beluga Whale", "Diverse vocalizations", "1,251 Hz", "6.19 bits", COLOR_RED),
        ]
        headers_text = ["Species", "Vocalization Type", "Centroid Freq.", "Entropy"]

    # Header row
    col_lefts = [Cm(1), Cm(8.5), Cm(17.5), Cm(24)]
    col_widths = [Cm(7), Cm(8.5), Cm(6), Cm(6)]
    for i, h in enumerate(headers_text):
        hbox = add_shape(slide, col_lefts[i], Cm(2.5), col_widths[i], Cm(1.2),
                         shape_type=MSO_SHAPE.RECTANGLE, fill_color=COLOR_TITLE_BG)
        set_text(hbox, h, font_size=Pt(12), font_color=COLOR_WHITE, bold=True)

    # Species rows
    for idx, (name, vtype, freq, entropy, color) in enumerate(species_info):
        y = Cm(4.0) + idx * Cm(2.2)

        # Color indicator + name
        indicator = add_shape(slide, Cm(1), y, Cm(0.5), Cm(1.5),
                              shape_type=MSO_SHAPE.RECTANGLE, fill_color=color)
        name_box = add_shape(slide, Cm(1.5), y, Cm(6.5), Cm(1.5),
                             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                             fill_color=COLOR_WHITE, line_color=color, line_width=Pt(1.5))
        set_text(name_box, name, font_size=Pt(12), font_color=color, bold=True)

        # Vocalization type
        vtype_box = add_shape(slide, Cm(8.5), y, Cm(8.5), Cm(1.5),
                              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                              fill_color=COLOR_WHITE, line_color=COLOR_GRAY, line_width=Pt(1))
        set_text(vtype_box, vtype, font_size=Pt(11), font_color=COLOR_BODY)

        # Frequency
        freq_box = add_shape(slide, Cm(17.5), y, Cm(6), Cm(1.5),
                             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                             fill_color=COLOR_WHITE, line_color=COLOR_GRAY, line_width=Pt(1))
        set_text(freq_box, freq, font_size=Pt(12), font_color=COLOR_BODY, bold=True)

        # Entropy (with bar)
        ent_val = float(entropy.split()[0])
        max_ent = 8.0
        bar_width = Cm(5.5 * ent_val / max_ent)
        add_shape(slide, Cm(24), y + Cm(0.15), bar_width, Cm(1.2),
                  shape_type=MSO_SHAPE.RECTANGLE, fill_color=color)
        ent_label = slide.shapes.add_textbox(Cm(24), y, Cm(6), Cm(1.5))
        etf = ent_label.text_frame
        ep = etf.paragraphs[0]
        ep.text = entropy
        ep.font.size = Pt(11)
        ep.font.color.rgb = COLOR_WHITE
        ep.font.bold = True
        ep.alignment = PP_ALIGN.CENTER

    # Footer
    note_text = "※ すべての要素は編集可能です" if lang == "ja" else "Note: All elements are editable"
    txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = note_text
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.italic = True

    return slide


def make_conclusion_diagram(prs, lang="ja"):
    """Create editable conclusion/summary diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    title_text = "結論と今後の展望" if lang == "ja" else "Conclusions and Future Directions"
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    if lang == "ja":
        findings = [
            ("種固有スペクトル符号", "6種は周波数特性で\n明確に区別可能", COLOR_BLUE, "支持"),
            ("CDMA的直交性", "種内相関 > 種間相関\np = 1.57×10⁻³", COLOR_GREEN, "支持"),
            ("非線形結合", "全種でバイコヒーレンス\n> 0を検出", COLOR_ORANGE, "部分的"),
            ("情報構造", "エントロピー 2.28〜7.32\n bits の種間差", COLOR_PURPLE, "支持"),
        ]
        future = [
            "複数ハイドロフォン同期録音による伝搬解析",
            "個体レベルの符号分離",
            "時系列的な符号変化の追跡",
            "深層学習による特徴抽出",
        ]
        future_title = "今後の課題"
    else:
        findings = [
            ("Species-specific codes", "6 species clearly\ndistinguishable by\nfrequency features", COLOR_BLUE, "Supported"),
            ("CDMA orthogonality", "Within > Between\ncorrelation\np = 1.57×10⁻³", COLOR_GREEN, "Supported"),
            ("Nonlinear coupling", "Bicoherence > 0\ndetected in all species", COLOR_ORANGE, "Partial"),
            ("Information structure", "Entropy 2.28-7.32 bits\ninterspecific variation", COLOR_PURPLE, "Supported"),
        ]
        future = [
            "Multi-hydrophone synchronized recordings",
            "Individual-level code separation",
            "Longitudinal code variation tracking",
            "Deep learning feature extraction",
        ]
        future_title = "Future Directions"

    # Findings boxes
    box_width = Cm(7.5)
    box_height = Cm(6.0)
    start_left = Cm(0.8)

    for i, (title, desc, color, status) in enumerate(findings):
        left = start_left + i * (box_width + Cm(0.5))
        top = Cm(2.8)

        # Main box
        box = add_shape(slide, left, top, box_width, box_height,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=COLOR_WHITE, line_color=color, line_width=Pt(2))

        # Title bar
        title_bar = add_shape(slide, left, top, box_width, Cm(1.3),
                              shape_type=MSO_SHAPE.RECTANGLE, fill_color=color)
        set_text(title_bar, title, font_size=Pt(11), font_color=COLOR_WHITE, bold=True)

        # Description
        desc_box = slide.shapes.add_textbox(left + Cm(0.3), top + Cm(1.5),
                                            box_width - Cm(0.6), Cm(3.0))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLOR_BODY
        dp.alignment = PP_ALIGN.CENTER

        # Status badge
        badge_color = COLOR_GREEN if status in ("支持", "Supported") else COLOR_ORANGE
        badge = add_shape(slide, left + Cm(1.5), top + box_height - Cm(1.5),
                          box_width - Cm(3), Cm(1.0),
                          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                          fill_color=badge_color)
        set_text(badge, status, font_size=Pt(10), font_color=COLOR_WHITE, bold=True)

    # Future directions section
    future_top = Cm(9.5)
    future_header = add_shape(slide, Cm(0.8), future_top, Cm(32), Cm(1.2),
                              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                              fill_color=COLOR_TITLE_BG)
    set_text(future_header, future_title, font_size=Pt(16), font_color=COLOR_WHITE, bold=True)

    for i, item in enumerate(future):
        left = Cm(0.8) + i * Cm(8)
        top = future_top + Cm(1.8)
        item_box = add_shape(slide, left, top, Cm(7.5), Cm(2.5),
                             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                             fill_color=COLOR_LIGHT_BG, line_color=COLOR_ACCENT1, line_width=Pt(1))
        tf_item = set_text(item_box, item, font_size=Pt(10), font_color=COLOR_ACCENT1)

        # Number circle
        num_circle = add_shape(slide, left + Cm(3), top - Cm(0.6), Cm(1.0), Cm(1.0),
                               shape_type=MSO_SHAPE.OVAL, fill_color=COLOR_ACCENT1)
        set_text(num_circle, str(i + 1), font_size=Pt(12), font_color=COLOR_WHITE, bold=True)

    # Footer
    note_text = "※ すべての要素は編集可能です" if lang == "ja" else "Note: All elements are editable"
    txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = note_text
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.italic = True

    return slide


def make_methodology_diagram(prs, lang="ja"):
    """Create editable methodology/analysis methods overview diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.3), fill_color=COLOR_TITLE_BG)

    title_text = "7つの解析手法" if lang == "ja" else "Seven Analytical Methods"
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_HEADING
    p.font.bold = True

    if lang == "ja":
        methods = [
            ("1", "スペクトログラム", "STFT → 周波数構造\nスペクトル重心・帯域幅", COLOR_BLUE),
            ("2", "ICI解析", "クリック間隔分布\n個体識別符号の検出", COLOR_GREEN),
            ("3", "バイスペクトル", "非線形位相結合\nうなり効果の検出", COLOR_ORANGE),
            ("4", "エントロピー", "Shannon情報量\nZipf則適合度", COLOR_PURPLE),
            ("5", "時間構造", "自己/相互相関\nスペクトル平坦度", COLOR_TEAL),
            ("6", "種間比較", "箱ひげ図\n特徴量の種間差", COLOR_RED),
            ("7", "CDMA直交性", "スペクトル相関行列\nMann-Whitney U検定", COLOR_ACCENT2),
        ]
    else:
        methods = [
            ("1", "Spectrogram", "STFT → frequency structure\ncentroid, bandwidth", COLOR_BLUE),
            ("2", "ICI Analysis", "Click interval distribution\nidentity code detection", COLOR_GREEN),
            ("3", "Bispectrum", "Nonlinear phase coupling\nbeat effect detection", COLOR_ORANGE),
            ("4", "Entropy", "Shannon information\nZipf's law fit", COLOR_PURPLE),
            ("5", "Temporal", "Auto/cross-correlation\nspectral flatness", COLOR_TEAL),
            ("6", "Cross-species", "Box plots\ninterspecific differences", COLOR_RED),
            ("7", "CDMA Orthogonality", "Spectral correlation matrix\nMann-Whitney U test", COLOR_ACCENT2),
        ]

    # Layout: 4 on top row, 3 on bottom row
    box_w = Cm(7.5)
    box_h = Cm(5.5)

    for i, (num, name, desc, color) in enumerate(methods):
        if i < 4:
            left = Cm(0.6) + i * (box_w + Cm(0.5))
            top = Cm(2.5)
        else:
            left = Cm(4.5) + (i - 4) * (box_w + Cm(0.5))
            top = Cm(8.5)

        # Box
        box = add_shape(slide, left, top, box_w, box_h,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=COLOR_WHITE, line_color=color, line_width=Pt(2))

        # Number circle
        num_circle = add_shape(slide, left + box_w / 2 - Cm(0.5), top - Cm(0.5),
                               Cm(1.0), Cm(1.0),
                               shape_type=MSO_SHAPE.OVAL, fill_color=color)
        set_text(num_circle, num, font_size=Pt(14), font_color=COLOR_WHITE, bold=True)

        # Name
        name_bar = add_shape(slide, left, top + Cm(0.6), box_w, Cm(1.2),
                             shape_type=MSO_SHAPE.RECTANGLE, fill_color=color)
        set_text(name_bar, name, font_size=Pt(12), font_color=COLOR_WHITE, bold=True)

        # Description
        desc_box = slide.shapes.add_textbox(left + Cm(0.3), top + Cm(2.0),
                                            box_w - Cm(0.6), box_h - Cm(2.5))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLOR_BODY
        dp.alignment = PP_ALIGN.CENTER

    # Footer
    note_text = "※ すべての要素は編集可能です" if lang == "ja" else "Note: All elements are editable"
    txBox2 = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(20), Cm(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = note_text
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.italic = True

    return slide


# ============================================================
# Main presentation builder
# ============================================================


def build_presentation(lang="ja"):
    """Build the full presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    species_names = SPECIES_JA if lang == "ja" else SPECIES_EN
    suffix = f"_{lang}.png"

    # ---- Title slide ----
    if lang == "ja":
        make_title_slide(prs,
                         "鯨類音響コミュニケーションにおける\nエンコーディング構造の定量的解析",
                         "CDMA的符号分割仮説およびうなり周波数仮説の検証")
    else:
        make_title_slide(prs,
                         "Quantitative Analysis of Encoding Structures\nin Cetacean Acoustic Communication",
                         "Verification of CDMA-like Code Division and Beat Frequency Hypotheses")

    # ---- Editable: Hypothesis diagram ----
    make_hypothesis_diagram(prs, lang)

    # ---- Editable: Species overview ----
    make_species_overview_diagram(prs, lang)

    # ---- Editable: Analysis pipeline flow ----
    make_analysis_pipeline_flow(prs, lang)

    # ---- Editable: Methodology diagram ----
    make_methodology_diagram(prs, lang)

    # ---- Table: Spectral features summary ----
    if lang == "ja":
        headers = ["種", "スペクトル重心 (Hz)", "帯域幅 (Hz)", "支配周波数 (Hz)", "エントロピー (bits)", "バイコヒーレンス"]
        rows = [
            ["マッコウクジラ", "2,309", "1,874", "1,047", "7.32", "0.000751"],
            ["ザトウクジラ", "593", "389", "508", "4.47", "0.000285"],
            ["シャチ", "1,804", "1,300", "1,547", "6.66", "0.000131"],
            ["ナガスクジラ", "38", "30", "29", "2.28", "0.000124"],
            ["バンドウイルカ", "4,021", "2,488", "1,430", "5.79", "0.000340"],
            ["シロイルカ", "1,251", "1,650", "128", "6.19", "0.000386"],
        ]
        caption = "表1. 6種の海獣の音響スペクトル特徴の比較"
        table_title = "スペクトル特徴の種間比較"
    else:
        headers = ["Species", "Centroid (Hz)", "Bandwidth (Hz)", "Dominant (Hz)", "Entropy (bits)", "Bicoherence"]
        rows = [
            ["Sperm Whale", "2,309", "1,874", "1,047", "7.32", "0.000751"],
            ["Humpback Whale", "593", "389", "508", "4.47", "0.000285"],
            ["Killer Whale", "1,804", "1,300", "1,547", "6.66", "0.000131"],
            ["Fin Whale", "38", "30", "29", "2.28", "0.000124"],
            ["Bottlenose Dolphin", "4,021", "2,488", "1,430", "5.79", "0.000340"],
            ["Beluga Whale", "1,251", "1,650", "128", "6.19", "0.000386"],
        ]
        caption = "Table 1. Comparison of acoustic spectral features across six marine mammal species"
        table_title = "Cross-Species Spectral Feature Comparison"
    make_table_slide(prs, table_title, headers, rows, caption)

    # ---- Table: CDMA orthogonality results ----
    if lang == "ja":
        headers2 = ["指標", "値"]
        rows2 = [
            ["種内スペクトル相関 (平均)", "0.4000 ± 0.2939"],
            ["種間スペクトル相関 (平均)", "0.2859 ± 0.2647"],
            ["Mann-Whitney U 統計量", "13,921.0"],
            ["p値", "1.57 × 10⁻³"],
            ["帰無仮説", "棄却 (α = 0.05)"],
            ["結論", "種固有スペクトル符号が存在"],
        ]
        caption2 = "表2. CDMA的直交性検定の結果"
        table_title2 = "CDMA的直交性検定結果"
    else:
        headers2 = ["Metric", "Value"]
        rows2 = [
            ["Within-species correlation (mean)", "0.4000 ± 0.2939"],
            ["Between-species correlation (mean)", "0.2859 ± 0.2647"],
            ["Mann-Whitney U statistic", "13,921.0"],
            ["p-value", "1.57 × 10⁻³"],
            ["Null hypothesis", "Rejected (α = 0.05)"],
            ["Conclusion", "Species-specific spectral codes exist"],
        ]
        caption2 = "Table 2. CDMA-like orthogonality test results"
        table_title2 = "CDMA-like Orthogonality Test Results"
    make_table_slide(prs, table_title2, headers2, rows2, caption2)

    # ---- Table: ICI summary ----
    if lang == "ja":
        headers3 = ["種", "サンプル", "中央ICI (ms)", "標準偏差 (ms)", "クリック数"]
        rows3 = [
            ["マッコウクジラ", "1", "19.7", "115.7", "25"],
            ["マッコウクジラ", "2", "18.5", "86.5", "16"],
            ["マッコウクジラ", "4", "119.1", "268.7", "208"],
            ["シャチ", "1", "10.2", "99.2", "25"],
            ["シャチ", "2", "19.6", "78.8", "30"],
            ["バンドウイルカ", "1", "4.7", "226.9", "6"],
            ["バンドウイルカ", "2", "10.4", "244.3", "12"],
        ]
        caption3 = "表3. クリック間隔 (ICI) 解析の要約"
        table_title3 = "クリック間隔 (ICI) 解析結果"
    else:
        headers3 = ["Species", "Sample", "Median ICI (ms)", "Std Dev (ms)", "Clicks"]
        rows3 = [
            ["Sperm Whale", "1", "19.7", "115.7", "25"],
            ["Sperm Whale", "2", "18.5", "86.5", "16"],
            ["Sperm Whale", "4", "119.1", "268.7", "208"],
            ["Killer Whale", "1", "10.2", "99.2", "25"],
            ["Killer Whale", "2", "19.6", "78.8", "30"],
            ["Bottlenose Dolphin", "1", "4.7", "226.9", "6"],
            ["Bottlenose Dolphin", "2", "10.4", "244.3", "12"],
        ]
        caption3 = "Table 3. Summary of inter-click interval (ICI) analysis"
        table_title3 = "Inter-Click Interval (ICI) Analysis Results"
    make_table_slide(prs, table_title3, headers3, rows3, caption3)

    # ---- Section: Spectrogram results ----
    section_title = "スペクトログラム解析結果" if lang == "ja" else "Spectrogram Analysis Results"
    make_section_slide(prs, section_title, section_num=1)

    for sp in TARGET_SPECIES:
        img_path = os.path.join(OUTPUT_DIR, f"spectrogram_{sp}{suffix}")
        if os.path.exists(img_path):
            sp_name = species_names.get(sp, sp)
            if lang == "ja":
                title = f"スペクトログラム: {sp_name}"
                caption = f"図. {sp_name}のスペクトログラムと周波数構造"
            else:
                title = f"Spectrogram: {sp_name}"
                caption = f"Figure. Spectrogram and frequency structure of {sp_name}"
            make_image_slide(prs, title, img_path, caption)

    # ---- Section: ICI results ----
    section_title = "クリック間隔 (ICI) 解析結果" if lang == "ja" else "Inter-Click Interval (ICI) Results"
    make_section_slide(prs, section_title, section_num=2)

    for sp in CLICK_SPECIES:
        img_path = os.path.join(OUTPUT_DIR, f"ici_{sp}{suffix}")
        if os.path.exists(img_path):
            sp_name = species_names.get(sp, sp)
            if lang == "ja":
                title = f"ICI分布: {sp_name}"
                caption = f"図. {sp_name}のクリック間隔分布"
            else:
                title = f"ICI Distribution: {sp_name}"
                caption = f"Figure. Inter-click interval distribution of {sp_name}"
            make_image_slide(prs, title, img_path, caption)

    # ---- Section: Bispectrum results ----
    section_title = "バイスペクトル解析結果" if lang == "ja" else "Bispectral Analysis Results"
    make_section_slide(prs, section_title, section_num=3)

    for sp in TARGET_SPECIES:
        img_path = os.path.join(OUTPUT_DIR, f"bispectrum_{sp}{suffix}")
        if os.path.exists(img_path):
            sp_name = species_names.get(sp, sp)
            if lang == "ja":
                title = f"バイスペクトル: {sp_name}"
                caption = f"図. {sp_name}のバイコヒーレンスマップ（非線形結合）"
            else:
                title = f"Bispectrum: {sp_name}"
                caption = f"Figure. Bicoherence map (nonlinear coupling) of {sp_name}"
            make_image_slide(prs, title, img_path, caption)

    # ---- Section: Entropy results ----
    section_title = "情報エントロピー解析結果" if lang == "ja" else "Information Entropy Results"
    make_section_slide(prs, section_title, section_num=4)

    for sp in TARGET_SPECIES:
        img_path = os.path.join(OUTPUT_DIR, f"entropy_{sp}{suffix}")
        if os.path.exists(img_path):
            sp_name = species_names.get(sp, sp)
            if lang == "ja":
                title = f"エントロピーとZipf分布: {sp_name}"
                caption = f"図. {sp_name}の情報エントロピーおよびZipf則適合度"
            else:
                title = f"Entropy & Zipf Distribution: {sp_name}"
                caption = f"Figure. Information entropy and Zipf's law fit of {sp_name}"
            make_image_slide(prs, title, img_path, caption)

    # ---- Section: Temporal structure ----
    section_title = "時間構造解析結果" if lang == "ja" else "Temporal Structure Results"
    make_section_slide(prs, section_title, section_num=5)

    for sp in TARGET_SPECIES:
        img_path = os.path.join(OUTPUT_DIR, f"temporal_{sp}{suffix}")
        if os.path.exists(img_path):
            sp_name = species_names.get(sp, sp)
            if lang == "ja":
                title = f"時間構造: {sp_name}"
                caption = f"図. {sp_name}の時間構造解析（自己相関、変調スペクトル等）"
            else:
                title = f"Temporal Structure: {sp_name}"
                caption = f"Figure. Temporal structure analysis of {sp_name}"
            make_image_slide(prs, title, img_path, caption)

    # ---- Section: Cross-species comparison ----
    section_title = "種間比較" if lang == "ja" else "Cross-Species Comparison"
    make_section_slide(prs, section_title, section_num=6)

    img_path = os.path.join(OUTPUT_DIR, f"cross_species{suffix}")
    if os.path.exists(img_path):
        if lang == "ja":
            make_image_slide(prs, "種間音響特徴比較", img_path,
                             "図. 6種の音響特徴の箱ひげ図比較")
        else:
            make_image_slide(prs, "Cross-Species Acoustic Feature Comparison", img_path,
                             "Figure. Box-plot comparison of acoustic features across six species")

    # ---- Section: CDMA orthogonality ----
    section_title = "CDMA的直交性解析" if lang == "ja" else "CDMA-like Orthogonality Analysis"
    make_section_slide(prs, section_title, section_num=7)

    img_path = os.path.join(OUTPUT_DIR, f"cdma{suffix}")
    if os.path.exists(img_path):
        if lang == "ja":
            make_image_slide(prs, "CDMA的直交性: スペクトル相関行列", img_path,
                             "図. 種内vs種間スペクトル相関の比較（Mann-Whitney U検定, p=1.57×10⁻³）")
        else:
            make_image_slide(prs, "CDMA Orthogonality: Spectral Correlation Matrix", img_path,
                             "Figure. Within vs between species spectral correlation (Mann-Whitney U, p=1.57×10⁻³)")

    # ---- Editable: Conclusion diagram ----
    make_conclusion_diagram(prs, lang)

    # Save
    if lang == "ja":
        fname = "鯨類音響コミュニケーション解析_図表集.pptx"
    else:
        fname = "Cetacean_Acoustic_Communication_Figures.pptx"
    path = os.path.join(PAPER_DIR, fname)
    prs.save(path)
    print(f"  Saved: {path} ({len(prs.slides)} slides)")
    return path


def main():
    print("=" * 70)
    print("Generating PowerPoint Presentations")
    print("pptxプレゼンテーション生成")
    print("=" * 70)

    print("\nGenerating Japanese presentation...")
    ja_path = build_presentation("ja")

    print("\nGenerating English presentation...")
    en_path = build_presentation("en")

    print(f"\n{'='*70}")
    print(f"Presentations saved:")
    print(f"  Japanese: {ja_path}")
    print(f"  English:  {en_path}")
    print(f"{'='*70}")


if __name__ == "__main__":
    main()
